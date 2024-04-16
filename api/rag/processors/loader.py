""" from https://github.com/chatchat-space/Langchain-Chatchat """
import csv
import os
from io import BytesIO, TextIOWrapper
from pathlib import Path
from typing import (
    Dict,
    List,
    Optional,
    Sequence,
    TYPE_CHECKING,
)

import nltk
import numpy as np
from PIL import Image
from langchain.docstore.document import Document
from langchain.document_loaders.helpers import detect_file_encodings
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.document_loaders.unstructured import UnstructuredFileLoader
from tqdm import tqdm

if TYPE_CHECKING:
    try:
        from rapidocr_paddle import RapidOCR
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR


PDF_OCR_THRESHOLD = (0.6, 0.6)

NLTK_DATA_PATH = os.path.join(Path(__file__).parents[3], "assets", "nltk_data")
nltk.data.path = [NLTK_DATA_PATH] + nltk.data.path


__all__ = [
    "FilteredCSVLoader",
    "RapidOCRDocLoader",
    "RapidOCRImageLoader",
    "RapidOCRPDFLoader",
    "OpenParserPDFLoader",
    "RapidOCRPPTLoader",
]


def get_ocr(use_cuda: bool = True) -> "RapidOCR":
    try:
        from rapidocr_paddle import RapidOCR
        ocr = RapidOCR(det_use_cuda=use_cuda, cls_use_cuda=use_cuda, rec_use_cuda=use_cuda)
    except ImportError:
        from rapidocr_onnxruntime import RapidOCR
        ocr = RapidOCR()
    return ocr


class FilteredCSVLoader(CSVLoader):
    def __init__(
        self,
        file_path: str,
        columns_to_read: List[str],
        source_column: Optional[str] = None,
        metadata_columns: Sequence[str] = (),
        csv_args: Optional[Dict] = None,
        encoding: Optional[str] = None,
        autodetect_encoding: bool = False,
    ):
        super().__init__(
            file_path=file_path,
            source_column=source_column,
            metadata_columns=metadata_columns,
            csv_args=csv_args,
            encoding=encoding,
            autodetect_encoding=autodetect_encoding,
        )
        self.columns_to_read = columns_to_read

    def load(self) -> List[Document]:
        """Load data into document objects."""
        docs = []
        try:
            with open(self.file_path, newline="", encoding=self.encoding) as csvfile:
                docs = self.__read_file(csvfile)
        except UnicodeDecodeError as e:
            if self.autodetect_encoding:
                detected_encodings = detect_file_encodings(self.file_path)
                for encoding in detected_encodings:
                    try:
                        with open(
                            self.file_path, newline="", encoding=encoding.encoding
                        ) as csvfile:
                            docs = self.__read_file(csvfile)
                            break
                    except UnicodeDecodeError:
                        continue
            else:
                raise RuntimeError(f"Error loading {self.file_path}") from e
        except Exception as e:
            raise RuntimeError(f"Error loading {self.file_path}") from e

        return docs

    def __read_file(self, csvfile: TextIOWrapper) -> List[Document]:
        docs = []
        csv_reader = csv.DictReader(csvfile, **self.csv_args)  # type: ignore
        for i, row in enumerate(csv_reader):
            content = []
            for col in self.columns_to_read:
                if col in row:
                    content.append(f'{col}:{str(row[col])}')
                else:
                    raise ValueError(f"Column '{self.columns_to_read[0]}' not found in CSV file.")
            content = '\n'.join(content)
            # Extract the source if available
            source = (
                row.get(self.source_column, None)
                if self.source_column is not None
                else self.file_path
            )
            metadata = {"source": source, "row": i}

            for col in self.metadata_columns:
                if col in row:
                    metadata[col] = row[col]

            doc = Document(page_content=content, metadata=metadata)
            docs.append(doc)

        return docs


class RapidOCRDocLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        from unstructured.partition.text import partition_text

        def doc2text(filepath):
            from docx.table import _Cell, Table
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.text.paragraph import Paragraph
            from docx import Document, ImagePart
            from PIL import Image

            ocr = get_ocr()
            doc = Document(filepath)
            resp = ""

            def iter_block_items(parent):
                from docx.document import Document
                if isinstance(parent, Document):
                    parent_elm = parent.element.body
                elif isinstance(parent, _Cell):
                    parent_elm = parent._tc
                else:
                    raise ValueError("RapidOCRDocLoader parse fail")

                for child in parent_elm.iterchildren():
                    if isinstance(child, CT_P):
                        yield Paragraph(child, parent)
                    elif isinstance(child, CT_Tbl):
                        yield Table(child, parent)

            b_unit = tqdm(
                total=len(doc.paragraphs)+len(doc.tables),
                desc="RapidOCRDocLoader block index: 0"
            )
            for i, block in enumerate(iter_block_items(doc)):
                b_unit.set_description(
                    "RapidOCRDocLoader  block index: {}".format(i))
                b_unit.refresh()
                if isinstance(block, Paragraph):
                    resp += block.text.strip() + "\n"
                    images = block._element.xpath('.//pic:pic')  # 获取所有图片
                    for image in images:
                        for img_id in image.xpath('.//a:blip/@r:embed'):  # 获取图片id
                            part = doc.part.related_parts[img_id]  # 根据图片id获取对应的图片
                            if isinstance(part, ImagePart):
                                image = Image.open(BytesIO(part._blob))
                                result, _ = ocr(np.array(image))
                                if result:
                                    ocr_result = [line[1] for line in result]
                                    resp += "\n".join(ocr_result)
                elif isinstance(block, Table):
                    for row in block.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                b_unit.update(1)
            return resp

        text = doc2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)


class RapidOCRImageLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        from unstructured.partition.text import partition_text

        def img2text(filepath):
            resp = ""
            ocr = get_ocr()
            result, _ = ocr(filepath)
            if result:
                ocr_result = [line[1] for line in result]
                resp += "\n".join(ocr_result)
            return resp

        text = img2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)


class RapidOCRPDFLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        import cv2
        from unstructured.partition.text import partition_text

        def rotate_img(img, angle):
            h, w = img.shape[:2]
            rotate_center = (w / 2, h / 2)
            # 获取旋转矩阵
            # 参数1为旋转中心点;
            # 参数2为旋转角度,正值-逆时针旋转;负值-顺时针旋转
            # 参数3为各向同性的比例因子,1.0原图，2.0变成原来的2倍，0.5变成原来的0.5倍
            M = cv2.getRotationMatrix2D(rotate_center, angle, 1.0)
            # 计算图像新边界
            new_w = int(h * np.abs(M[0, 1]) + w * np.abs(M[0, 0]))
            new_h = int(h * np.abs(M[0, 0]) + w * np.abs(M[0, 1]))
            # 调整旋转矩阵以考虑平移
            M[0, 2] += (new_w - w) / 2
            M[1, 2] += (new_h - h) / 2

            rotated_img = cv2.warpAffine(img, M, (new_w, new_h))
            return rotated_img

        def pdf2text(filepath):
            import fitz  # pyMuPDF里面的fitz包，不要与pip install fitz混淆

            ocr = get_ocr()
            doc = fitz.open(filepath)
            resp = ""

            b_unit = tqdm(total=doc.page_count, desc="RapidOCRPDFLoader context page index: 0")
            for i, page in enumerate(doc):
                b_unit.set_description("RapidOCRPDFLoader context page index: {}".format(i))
                b_unit.refresh()
                text = page.get_text("")
                resp += text + "\n"

                img_list = page.get_image_info(xrefs=True)
                for img in img_list:
                    if xref := img.get("xref"):
                        bbox = img["bbox"]
                        # 检查图片尺寸是否超过设定的阈值
                        if ((bbox[2] - bbox[0]) / page.rect.width < PDF_OCR_THRESHOLD[0]
                                or (bbox[3] - bbox[1]) / page.rect.height < PDF_OCR_THRESHOLD[1]):
                            continue
                        pix = fitz.Pixmap(doc, xref)
                        if int(page.rotation) != 0:  # 如果Page有旋转角度，则旋转图片
                            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)
                            tmp_img = Image.fromarray(img_array)
                            ori_img = cv2.cvtColor(np.array(tmp_img), cv2.COLOR_RGB2BGR)
                            rot_img = rotate_img(img=ori_img, angle=360 - page.rotation)
                            img_array = cv2.cvtColor(rot_img, cv2.COLOR_RGB2BGR)
                        else:
                            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, -1)

                        result, _ = ocr(img_array)
                        if result:
                            ocr_result = [line[1] for line in result]
                            resp += "\n".join(ocr_result)

                # 更新进度
                b_unit.update(1)
            return resp

        text = pdf2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)


class OpenParserPDFLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        from unstructured.partition.text import partition_text

        def pdf2text(filepath):
            from openparse import DocumentParser

            parser = DocumentParser()
            parsed_content = parser.parse(filepath)
            resp = "\n".join(node.text for node in parsed_content.nodes)
            return resp

        text = pdf2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)


class RapidOCRPPTLoader(UnstructuredFileLoader):
    def _get_elements(self) -> List:
        from unstructured.partition.text import partition_text

        def ppt2text(filepath):
            from pptx import Presentation

            ocr = get_ocr()
            prs = Presentation(filepath)
            resp = ""

            def extract_text(shape):
                nonlocal resp
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            b_unit = tqdm(total=len(prs.slides), desc="RapidOCRPPTLoader slide index: 1")
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number)
                )
                b_unit.refresh()
                sorted_shapes = sorted(
                    slide.shapes, key=lambda x: (x.top, x.left)
                )  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)
